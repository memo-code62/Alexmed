import streamlit as st
import os
import io
from pathlib import Path
from PIL import Image

# الاستيرادات الأساسية لـ LlamaIndex
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, StorageContext, load_index_from_storage
from llama_index.readers.web import SimpleWebPageReader
from llama_index.embeddings.gemini import GeminiEmbedding
# استيرادات Gemini (تم التأكد من استخدام Gemini بدلاً من GeminiMultiModal)
from llama_index.llms.gemini import Gemini

# استيراد حزمة جوجل لتوليد الصور
from google import genai 

# استيراد مكتبات التنزيل
from docx import Document
from pptx import Presentation
from pptx.util import Inches 

# ---------------------------------
# 1. الإعدادات الأساسية والثوابت
# ---------------------------------

# **هام:** ضع مفتاحك في Secrets (GEMINI_API_KEY) 
INDEX_STORAGE_DIR = "storage"
PDF_DIR = "./" 
IMAGE_GENERATION_MODEL = 'imagen-3.0-generate-002' 

MEDICAL_URLS = [
    "https://pubmed.ncbi.nlm.nih.gov/", 
    "https://www.who.int/ar", 
    "https://www.cdc.gov/",
    "https://www.mayoclinic.org/",
    "https://www.medscape.com/",
    "https://www.hopkinsmedicine.org/",
]

SYSTEM_PROMPT = (
    "أنت مساعد طبي ذكي متخصص في الإجابة على استفسارات طلاب الطب. "
    "يجب عليك تحليل النص والصورة المرفقة (إن وجدت) واستخدامها مع المراجع المسترجعة للإجابة. "
    "الإجابة يجب أن تكون باللغة العربية، مع الحفاظ على **المصطلحات الطبية الأساسية (الأمراض، الأدوية، المصطلحات التشريحية)** باللغة الإنجليزية/اللاتينية داخل الأقواس. "
    "يجب أن تقدم إجاباتك في شكل منظم، وتستخدم الجداول والعناصر المرقمة عند طلب المقارنات. كما يمكنك إنشاء أسئلة تدريبية وتلخيصات و Mnemonic Devices (تحشيشات) عند طلبها."
)
# ---------------------------------
# 2. بناء/تحميل الفهرس المتعدد الأنماط
# ---------------------------------

@st.cache_resource
def setup_rag_engine():
    
    if "GEMINI_API_KEY" not in os.environ and not Path(INDEX_STORAGE_DIR).exists():
        st.error("❌ المفتاح السري لـ Gemini مفقود! يرجى إضافته في Secrets.")
        return None

    try:
        # 1. تهيئة نماذج اللغة (LLMs)
        llm_multi = Gemini(model="gemini-2.5-flash")
        llm_text = Gemini(model="gemini-2.5-flash")
        
        # 2. تهيئة نموذج التضمين (Embedding Model) لحل مشكلة OpenAI
        embed_model = GeminiEmbedding(model_name="text-embedding-004") 
        
    except Exception as e:
        st.error(f"❌ فشل تهيئة نموذج Gemini: {e}")
        return None

    if Path(INDEX_STORAGE_DIR).exists():
        st.info("🔄 جاري تحميل قاعدة المعرفة المتعددة الأنماط الموجودة مسبقًا...")
        storage_context = StorageContext.from_defaults(persist_dir=INDEX_STORAGE_DIR)
        index = load_index_from_storage(storage_context, llm=llm_text)
        
    else:
        st.warning("⏳ جاري بناء قاعدة المعرفة المتعددة الأنماط (قد يستغرق وقتًا طويلاً)...")
        
        try:
            # 1. قراءة الملفات المحلية (تم حذف الوسيطات غير المدعومة)
            pdf_documents = SimpleDirectoryReader(
                input_dir=PDF_DIR, 
                required_exts=[".pdf", ".jpg", ".png"]
            ).load_data()
            
            # 2. قراءة المواقع الإلكترونية (تم حل مشكلة التوافق)
            url_documents = SimpleWebPageReader().load_data(urls=MEDICAL_URLS)
            
            documents = pdf_documents + url_documents
            st.info(f"تم تحميل {len(documents)} مستند (نصي وبصري). جاري الفهرسة...")

            index = VectorStoreIndex.from_documents(
                documents,
                llm=llm_multi,
                embed_model=embed_model, # تمرير نموذج التضمين الخاص بـ Gemini
            )
            index.storage_context.persist(persist_dir=INDEX_STORAGE_DIR)
            st.success("✅ تم بناء قاعدة المعرفة المتعددة الأنماط وحفظها بنجاح! التطبيق جاهز.")
            
        except Exception as e:
            st.error(f"❌ خطأ حرج في بناء الفهرس: {e}")
            return None

    query_engine = index.as_query_engine(
        llm=llm_text,
        system_prompt=SYSTEM_PROMPT,
        streaming=True
    )
    return query_engine
    
# ---------------------------------
# 3. دوال توليد الوسائط والتنزيل
# ---------------------------------

@st.cache_resource
def get_gemini_client():
    return genai.Client()

def generate_image(prompt: str, is_animation=False):
    """تستخدم نموذج Imagen لإنشاء صورة ثابتة (جرافيك) بناءً على الوصف."""
    client = get_gemini_client()

    full_prompt = (
        f"Detailed medical diagram, high-quality colorful graphic design, "
        f"featuring arrows and explanatory labels showing the {prompt}"
    )

    st.info(f"🎨 جاري توليد جرافيك توضيحي لـ: {prompt}")

    try:
        if is_animation:
            st.warning("توليد الفيديو/الصور المتحركة معقد في Streamlit. سيتم توليد صورة ثابتة بدلاً من ذلك.")

        result = client.models.generate_images(
            model=IMAGE_GENERATION_MODEL,
            prompt=full_prompt,
            config=dict(
                number_of_images=1,
                output_mime_type="image/jpeg",
                aspect_ratio="16:9"
            )
        )

        if result.generated_images:
            image_data = result.generated_images[0].image.image_bytes
            return image_data, None
        else:
            return None, "لم يتمكن النموذج من توليد صورة بناءً على الوصف."

    except Exception as e:
        if "API_KEY_INVALID" in str(e):
            return None, "خطأ: مفتاح Gemini API غير صالح أو غير مهيأ لخدمة Imagen."
        return None, f"خطأ في توليد الصورة: {e}"

def convert_text_to_docx(text_content):
    """تحول النص إلى ملف Word (DOCX)."""
    document = Document()
    document.add_paragraph(text_content)
    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def convert_text_to_pptx(text_content):
    """تحول النص إلى عرض تقديمي (PPTX) بتقسيم بسيط."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] 

    paragraphs = text_content.split('\n\n') 

    for i, p in enumerate(paragraphs):
        if not p.strip(): continue

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"شريحة {i+1}: " + p.split('\n')[0][:50] + "..."

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = p

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def handle_image_generation(content, is_animation=False):
    """وظيفة مساعدة لمعالجة طلب توليد الصورة وعرضها."""
    image_prompt = content[:200]
    image_bytes, error = generate_image(image_prompt, is_animation=is_animation)

    if image_bytes:
        st.image(image_bytes, caption=f"صورة توضيحية تم توليدها لـ: {image_prompt}...")

        st.download_button(
            label="⬇️ تنزيل الصورة كـ JPG",
            data=image_bytes,
            file_name="medical_graphic.jpg",
            mime="image/jpeg",
            key=f"download_img_{hash(content)}"
        )
    else:
        st.error(f"فشل التوليد: {error}")

# ---------------------------------
# 4. واجهة Streamlit (التطبيق الرئيسي)
# ---------------------------------

st.set_page_config(page_title="مساعدك الطبي (RAG+Vision)", layout="wide")
st.title("👁️ مساعدك الطبي البصري (RAG+Vision)")
st.caption("يحلل ملفاتك وصورك المرفوعة للإجابة، ويدعم التنزيل المتعدد.")

query_engine = setup_rag_engine()

if query_engine:
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # عرض الرسائل السابقة
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

            # أزرار التوليد والتنزيل لرد المساعد
            if message["role"] == "assistant":
                col1, col2, col3, col4 = st.columns(4)

                with col1:
                    if st.button("🖼️ توليد جرافيك", key=f"img_{hash(message['content'])}"):
                        handle_image_generation(message['content'], is_animation=False)

                with col2:
                    docx_file = convert_text_to_docx(message["content"])
                    st.download_button(
                        label="📄 تنزيل كـ Word",
                        data=docx_file,
                        file_name="summary.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=f"download_docx_{hash(message['content'])}" 
                    )

                with col3:
                    pptx_file = convert_text_to_pptx(message["content"])
                    st.download_button(
                        label="📊 تنزيل كـ PPTX",
                        data=pptx_file,
                        file_name="presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_pptx_{hash(message['content'])}" 
                    )

                with col4:
                    st.download_button(
                        label="📜 تنزيل كـ TXT",
                        data=message["content"],
                        file_name="summary.txt",
                        mime="text/plain",
                        key=f"download_txt_{hash(message['content'])}" 
                    )

    # >>> منطقة تحميل الصورة والسؤال
    uploaded_file = st.file_uploader("🖼️ ارفع صورة طبية للسؤال عنها (اختياري)", type=["png", "jpg", "jpeg"])

    if prompt := st.chat_input("اطرح سؤالاً طبياً، ويمكنك إرفاق صورة..."):

        user_message = {"role": "user", "content": prompt}
        text_and_image_input = prompt
        image_to_query = None

        if uploaded_file:
            image = Image.open(uploaded_file)
            st.image(image, caption="الصورة المرفوعة", width=200)

            image_to_query = [image] 
            text_and_image_input = f"حلل الصورة المرفوعة معتمداً على مراجعك، ثم أجب عن السؤال: {prompt}"

        st.session_state.messages.append(user_message)
        with st.chat_message("user"):
            st.markdown(text_and_image_input)

        with st.chat_message("assistant"):
            response = query_engine.query(text_and_image_input, images=image_to_query) 

            st.write_stream(response.response_gen)
            st.session_state.messages.append({"role": "assistant", "content": response.response})